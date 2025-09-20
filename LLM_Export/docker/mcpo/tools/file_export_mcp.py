import re
import os
import ast
import json
import uuid
import emoji
import time
import base64
import datetime
import tarfile
import zipfile
import py7zr
import logging
import requests
from requests.auth import HTTPBasicAuth
import threading
import markdown2
import tempfile
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.shared import qn
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from docx.shared import Pt as DocxPt
from bs4 import BeautifulSoup, NavigableString
from mcp.server.fastmcp import FastMCP
from openpyxl import Workbook
import csv
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image as ReportLabImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm


PERSISTENT_FILES = os.getenv("PERSISTENT_FILES", "false")
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)


BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)


DOCS_TEMPLATE_PATH = os.getenv("DOCS_TEMPLATE_DIR", "/rootPath/templates")
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    # Search for .pptx, .docx, .xlsx templates inside DOCS_TEMPLATE_PATH
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
        logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    if XLSX_TEMPLATE_PATH:
        XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
        logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")




def search_image(query):
    log.debug(f"Searching for image with query: '{query}'")
    image_source = os.getenv("IMAGE_SOURCE", "unsplash")

    if image_source == "unsplash":
        return search_unsplash(query)
    elif image_source == "local_sd":
        return search_local_sd(query)
    elif image_source == "pexels":
        return search_pexels(query)
    else:
        log.warning(f"Image source unknown : {image_source}")
        return None

def search_local_sd(query: str):
    log.debug(f"Searching for local SD image with query: '{query}'")
    SD_URL = os.getenv("LOCAL_SD_URL")
    SD_USERNAME = os.getenv("LOCAL_SD_USERNAME")
    SD_PASSWORD = os.getenv("LOCAL_SD_PASSWORD")
    DEFAULT_MODEL = os.getenv("LOCAL_SD_DEFAULT_MODEL", "sd_xl_base_1.0.safetensors")
    DEFAULT_STEPS = int(os.getenv("LOCAL_SD_STEPS", 20))
    DEFAULT_WIDTH = int(os.getenv("LOCAL_SD_WIDTH", 512))
    DEFAULT_HEIGHT = int(os.getenv("LOCAL_SD_HEIGHT", 512))
    DEFAULT_CFG_SCALE = float(os.getenv("LOCAL_SD_CFG_SCALE", 1.5))
    DEFAULT_SCHEDULER = os.getenv("LOCAL_SD_SCHEDULER", "Karras")
    DEFAULT_SAMPLE = os.getenv("LOCAL_SD_SAMPLE", "Euler a")

    if not SD_URL:
        log.warning("LOCAL_SD_URL is not defined.")
        return None

    payload = {
        "prompt": query.strip(),
        "steps": DEFAULT_STEPS,
        "width": DEFAULT_WIDTH,
        "height": DEFAULT_HEIGHT,
        "cfg_scale": DEFAULT_CFG_SCALE,
        "sampler_name": DEFAULT_SAMPLE,
        "scheduler": DEFAULT_SCHEDULER,
        "enable_hr": False,
        "hr_upscaler": "Latent",
        "seed": -1,
        "override_settings": {
            "sd_model_checkpoint": DEFAULT_MODEL
        }
    }

    try:
        url = f"{SD_URL}/sdapi/v1/txt2img"
        log.debug(f"Sending request to local SD API at {url}")
        response = requests.post(
            url,
            json=payload,
            headers={"Content-Type": "application/json"},
            auth=HTTPBasicAuth(SD_USERNAME, SD_PASSWORD),
            timeout=30
        )
        response.raise_for_status()
        data = response.json()

        images = data.get("images", [])
        if not images:
            log.warning(f"No image generated for the request : '{query}'")
            return None

        image_b64 = images[0]
        image_data = base64.b64decode(image_b64)

        folder_path = _generate_unique_folder()
        filename = f"{query.replace(' ', '_')}.png"
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)

        with open(filepath, "wb") as f:
            f.write(image_data)

        return _public_url(folder_path, filename)

    except requests.exceptions.Timeout:
        log.error(f"Timeout during generation for : '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error : {e}")
    except Exception as e:
        log.error(f"Unexpected error : {e}")

    return None

def search_unsplash(query):
    log.debug(f"Searching Unsplash for query: '{query}'")
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        log.warning("UNSPLASH_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"Client-ID {api_key}"}
    log.debug(f"Sending request to Unsplash API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Unsplash API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Unsplash for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Unsplash for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None 

def search_pexels(query):
    log.debug(f"Searching Pexels for query: '{query}'")
    api_key = os.getenv("PEXELS_ACCESS_KEY")
    if not api_key:
        log.warning("PEXELS_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.pexels.com/v1/search"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"{api_key}"}
    log.debug(f"Sending request to Pexels API")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Pexels API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("photos"):
            image_url = data["photos"][0]["src"]["large"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Pexels for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Pexels for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None

def _resolve_log_level(val: str | None) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(
    level=_resolve_log_level(LOG_LEVEL_ENV),
    format=LOG_FORMAT_ENV,
)
log = logging.getLogger("file_export_mcp")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

mcp = FastMCP("file_export")

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    else:
        new_size = int(base_size / ratio)
        return PptPt(max(min_size, new_size))

def _public_url(folder_path: str, filename: str) -> str:
    """Build a stable public URL for a generated file."""
    folder = os.path.basename(folder_path).lstrip("/")
    name = filename.lstrip("/")
    return f"{BASE_URL}/{folder}/{name}"

def _generate_unique_folder() -> str:
    folder_name = f"export_{uuid.uuid4().hex[:10]}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = os.path.join(EXPORT_DIR, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def _generate_filename(folder_path: str, ext: str, filename: str = None) -> tuple[str, str]:
    if not filename:
        filename = f"export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    return filepath, filename

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(
    name="CustomHeading1",
    parent=styles["Heading1"],
    textColor=colors.HexColor("#0A1F44"),
    fontSize=18,
    spaceAfter=16,
    spaceBefore=12,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading2",
    parent=styles["Heading2"],
    textColor=colors.HexColor("#1C3F77"),
    fontSize=14,
    spaceAfter=12,
    spaceBefore=10,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading3",
    parent=styles["Heading3"],
    textColor=colors.HexColor("#3A6FB0"), 
    fontSize=12,
    spaceAfter=10,
    spaceBefore=8,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomNormal",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomListItem",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomCode",
    parent=styles["Code"],
    fontSize=10,
    leading=12,
    fontName="Courier",
    backColor=colors.HexColor("#F5F5F5"),
    borderColor=colors.HexColor("#CCCCCC"),
    borderWidth=1,
    leftIndent=10,
    rightIndent=10,
    topPadding=5,
    bottomPadding=5
))

def render_text_with_emojis(text: str) -> str:
    if not text:
        return ""
    try:
        converted = emoji.emojize(text, language="alias")
        return converted
    except Exception as e:
        log.error(f"Error in emoji conversion: {e}")
        return text

def process_list_items(ul_or_ol_element, is_ordered=False):
    items = []
    bullet_type = '1' if is_ordered else 'bullet'
    for li in ul_or_ol_element.find_all('li', recursive=False):
        li_text_parts = []
        for content in li.contents:
            if isinstance(content, NavigableString):
                li_text_parts.append(str(content))
            elif content.name not in ['ul', 'ol']:
                 li_text_parts.append(content.get_text())
        li_text = ''.join(li_text_parts).strip()
        list_item_paragraph = None
        if li_text:
            rendered_text = render_text_with_emojis(li_text)
            list_item_paragraph = Paragraph(rendered_text, styles["CustomListItem"])
        sub_lists = li.find_all(['ul', 'ol'], recursive=False)
        sub_flowables = []
        if list_item_paragraph:
             sub_flowables.append(list_item_paragraph)
        for sub_list in sub_lists:
            is_sub_ordered = sub_list.name == 'ol'
            nested_items = process_list_items(sub_list, is_sub_ordered)
            if nested_items:
                nested_list_flowable = ListFlowable(
                    nested_items,
                    bulletType='1' if is_sub_ordered else 'bullet',
                    leftIndent=10 * mm,
                    bulletIndent=5 * mm,
                    spaceBefore=2,
                    spaceAfter=2
                )
                sub_flowables.append(nested_list_flowable)
        if sub_flowables:
            items.append(ListItem(sub_flowables))
    return items

def render_html_elements(soup):
    log.debug("Starting render_html_elements...")
    story = []
    element_count = 0
    for elem in soup.children:
        element_count += 1
        log.debug(f"Processing element #{element_count}: {type(elem)}, name={getattr(elem, 'name', 'NavigableString')}")
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text:
                log.debug(f"Adding Paragraph from NavigableString: {text[:50]}...")
                story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                story.append(Spacer(1, 6))
        elif hasattr(elem, 'name'):
            tag_name = elem.name
            log.debug(f"Handling tag: <{tag_name}>")
            if tag_name == "h1":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H1: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading1"]))
                story.append(Spacer(1, 10))
            elif tag_name == "h2":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H2: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading2"]))
                story.append(Spacer(1, 8))
            elif tag_name == "h3":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H3: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading3"]))
                story.append(Spacer(1, 6))
            elif tag_name == "p":
                imgs = elem.find_all("img")
                if imgs:
                    for img_tag in imgs:
                        src = img_tag.get("src")
                        alt = img_tag.get("alt", "[Image]")
                        try:
                            if src and src.startswith("http"):
                                log.debug(f"Downloading image from URL: {src}")
                                response = requests.get(src)
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = Image(img_data, width=200, height=150)
                            else:
                                log.debug(f"Loading local image: {src}")
                                img = Image(src, width=200, height=150)
                            story.append(img)
                            story.append(Spacer(1, 10))
                        except Exception as e:
                            log.error(f"Error loading image {src}: {e}")
                            story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                            story.append(Spacer(1, 6))
                else:
                    text = render_text_with_emojis(elem.get_text().strip())
                    if text:
                        log.debug(f"Adding Paragraph: {text[:50]}...")
                        story.append(Paragraph(text, styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
            elif tag_name in ["ul", "ol"]:
                is_ordered = tag_name == "ol"
                log.debug(f"Processing list (ordered={is_ordered})...")
                items = process_list_items(elem, is_ordered)
                if items:
                    log.debug(f"Adding ListFlowable with {len(items)} items")
                    story.append(ListFlowable(items,
                        bulletType='1' if is_ordered else 'bullet',
                        leftIndent=10 * mm,
                        bulletIndent=5 * mm,
                        spaceBefore=6,
                        spaceAfter=10
                    ))
            elif tag_name == "blockquote":
                text = render_text_with_emojis(elem.get_text().strip())
                if text:
                    log.debug(f"Adding Blockquote: {text[:50]}...")
                    story.append(Paragraph(f"{text}", styles["CustomNormal"]))
                    story.append(Spacer(1, 8))
            elif tag_name in ["code", "pre"]:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Code/Pre block: {text[:50]}...")
                    story.append(Paragraph(text, styles["CustomCode"]))
                    story.append(Spacer(1, 6 if tag_name == "code" else 8))
            elif tag_name == "img":
                src = elem.get("src")
                alt = elem.get("alt", "[Image]")
                log.debug(f"Found <img> tag. src='{src}', alt='{alt}'")
                if src is not None: 
                    try:
                        if src.startswith("image_query:"):

                            query = src.replace("image_query:", "").strip()
                            log.debug(f"Handling image_query: '{query}'")
                            image_url = search_image(query)
                            if image_url:
                                log.debug(f"Downloading image from Unsplash URL: {image_url}")
                                response = requests.get(image_url)
                                log.debug(f"Image download response status: {response.status_code}")
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Unsplash)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                                log.warning(f"No image found for query: {query}")
                                story.append(Paragraph(f"[Image non trouvee pour: {query}]", styles["CustomNormal"]))
                                story.append(Spacer(1, 6))
                        elif src.startswith("http"):
                            log.debug(f"Downloading image from direct URL: {src}")
                            response = requests.get(src)
                            log.debug(f"Image download response status: {response.status_code}")
                            response.raise_for_status()
                            img_data = BytesIO(response.content)
                            img = ReportLabImage(img_data, width=200, height=150)
                            log.debug("Adding ReportLab Image object to story (Direct URL)")
                            story.append(img)
                            story.append(Spacer(1, 10))
                        else:
                            log.debug(f"Loading local image: {src}")
                            if os.path.exists(src):
                                img = ReportLabImage(src, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Local)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                               log.error(f"Local image file not found: {src}")
                               story.append(Paragraph(f"[Image locale non trouvee: {src}]", styles["CustomNormal"]))
                               story.append(Spacer(1, 6))
                    except requests.exceptions.RequestException as e:
                        log.error(f"Network error loading image {src}: {e}")
                        story.append(Paragraph(f"[Image (erreur reseau): {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        log.error(f"Error processing image {src}: {e}", exc_info=True) 
                        story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                else:
                    log.warning("Image tag found with no 'src' attribute.")
                    story.append(Paragraph(f"[Image: {alt} (source manquante)]", styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
            elif tag_name == "br":
                log.debug("Adding Spacer for <br>")
                story.append(Spacer(1, 6))
            else:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Paragraph for unknown tag <{tag_name}>: {text[:50]}...")
                    story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
    log.debug(f"Finished render_html_elements. Story contains {len(story)} elements.")
    return story

def _cleanup_files(folder_path: str, delay_minutes: int):
    def delete_files():
        time.sleep(delay_minutes * 60)
        try:
            import shutil
            shutil.rmtree(folder_path) 
            log.debug(f"Folder {folder_path} deleted.")
        except Exception as e:
            logging.error(f"Error deleting files : {e}")
    thread = threading.Thread(target=delete_files)
    thread.start()

def _convert_markdown_to_structured(markdown_content):
    """
    Converts Markdown content into a structured format for Word
    
    Args:
        markdown_content (str): Markdown content
        
    Returns:
        list: List of objects with 'text' and 'type'
    """
    if not markdown_content or not isinstance(markdown_content, str):
        return []
    
    lines = markdown_content.split('\n')
    structured = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('# '):
            structured.append({"text": line[2:].strip(), "type": "title"})
        elif line.startswith('## '):
            structured.append({"text": line[3:].strip(), "type": "heading"})
        elif line.startswith('### '):
            structured.append({"text": line[4:].strip(), "type": "subheading"})
        elif line.startswith('#### '):
            structured.append({"text": line[5:].strip(), "type": "subheading"})
        elif line.startswith('- '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('* '):
            structured.append({"text": line[2:].strip(), "type": "bullet"})
        elif line.startswith('**') and line.endswith('**'):
            structured.append({"text": line[2:-2].strip(), "type": "bold"})
        else:
            structured.append({"text": line, "type": "paragraph"})
    
    return structured

def _create_excel(data: list[list[str]], filename: str, folder_path: str | None = None) -> dict:
    log.debug("Creating Excel file")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "xlsx")

    wb = Workbook()
    ws = wb.active
    if isinstance(data, list):
        for row in data:
            ws.append(row)
    wb.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_csv(data: list[list[str]], filename: str, folder_path: str | None = None) -> dict:
    log.debug("Creating CSV file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "csv")

    with open(filepath, "w", newline="", encoding="utf-8") as f:
        if isinstance(data, list):
            csv.writer(f).writerows(data)
        else:
            csv.writer(f).writerow([data])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_pdf(text: str | list[str], filename: str, folder_path: str | None = None) -> dict:    
    log.debug("Creating PDF file")
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pdf")

    md_parts = []
    if isinstance(text, list):
        for item in text:
            if isinstance(item, str):
                md_parts.append(item)
            elif isinstance(item, dict):
                t = item.get("type")
                if t == "title":
                    md_parts.append(f"# {item.get('text','')}")
                elif t == "subtitle":
                    md_parts.append(f"## {item.get('text','')}")
                elif t == "paragraph":
                    md_parts.append(item.get("text",""))
                elif t == "list":
                    md_parts.append("\n".join([f"- {x}" for x in item.get("items",[])]))
                elif t in ("image","image_query"):
                    query = item.get("query","")
                    if query:
                        md_parts.append(f"![Image](image_query: {query})")
    else:
        md_parts = [str(text or "")]
        
    md_text = "\n\n".join(md_parts)    
   
    def replace_image_query(match):
        query = match.group(1).strip()
        image_url = search_image(query)
        return f'\n\n<img src="{image_url}" alt="Image: {query}" />\n\n' if image_url else ""

    md_text = re.sub(r'!\[[^\]]*\]\(\s*image_query:\s*([^)]+)\)', replace_image_query, md_text)
    html = markdown2.markdown(md_text, extras=['fenced-code-blocks','tables','break-on-newline','cuddled-lists'])
    soup = BeautifulSoup(html, "html.parser")
    story = render_html_elements(soup) or [Paragraph("Empty Content", styles["CustomNormal"])]

    doc = SimpleDocTemplate(filepath, topMargin=72, bottomMargin=72, leftMargin=72, rightMargin=72)
    try:
        doc.build(story)
    except Exception as e:
        log.error(f"Error building PDF {fname}: {e}", exc_info=True)
        doc2 = SimpleDocTemplate(filepath)
        doc2.build([Paragraph("Error in PDF generation", styles["CustomNormal"])])

    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_presentation(slides_data: list[dict], filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pptx")
      
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            log.debug("Attempting to load template...")
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                log.debug("Template is a Presentation object, converting to BytesIO")
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf); buf.seek(0)
                src = buf

            tmp = Presentation(src)
            log.debug(f"Template loaded with {len(tmp.slides)} slides")
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                # If the template has 2+ slides: slide 0 = title layout, slide 1 = content layout
                # If it has exactly 1 slide: use slide 0 layout for BOTH title and content
                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                log.debug("Using template layouts")

                # Keep only the first slide (as title base)
                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId  # type: ignore[attr-defined]
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]        # type: ignore[attr-defined]
            # else -> fall through to no-template
        except Exception:
            log.error(f"Error loading template: {e}")
            use_template = False
            prs = None

    if not use_template:
        log.debug("No valid template, creating new presentation with default layouts")
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    # Title slide (either existing template title, or newly added)
    if use_template:
        log.debug("Using template title slide")
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True
    else:
        log.debug("Creating new title slide")
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True

    # slide size in inches (for robust layout math)
    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    log.debug(f"Slide dimensions: {slide_w_in} x {slide_h_in} inches")

    # shared margins/gutters
    page_margin = 0.5   # outer margin on each side (inches)
    gutter = 0.3        # space between image and text (inches)

    # --- shared path: add content slides ---
    for i, slide_data in enumerate(slides_data):
        log.debug(f"Processing slide {i+1}: {slide_data.get('title', 'Untitled')}")
        if not isinstance(slide_data, dict):
            log.warning(f"Slide data is not a dict, skipping slide {i+1}")
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]
        log.debug(f"Adding slide with title: '{slide_title}'")
        slide = prs.slides.add_slide(content_layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28); r.font.bold = True

        # Find or create a content shape
        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph; break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph; break
                    except Exception:
                        pass
        except Exception:
            log.error(f"Error finding content placeholder: {e}")
            pass

        # Calculate title bottom position for proper image placement
        title_bottom_in = 1.0  # default fallback
        if slide.shapes.title:
            try:
                # Convert EMU to inches for title bottom position
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0)
                # Add small padding below title
                title_bottom_in += 0.2
            except Exception:
                title_bottom_in = 1.2  # fallback with padding

        if content_shape is None:

            content_shape = slide.shapes.add_textbox(Inches(page_margin), Inches(title_bottom_in), Inches(slide_w_in - 2*page_margin), Inches(slide_h_in - title_bottom_in - page_margin))
            log.debug("Creating new textbox for content")
        # prep text frame: wrap + shrink-to-fit + small inner margins
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            log.error(f"Error clearing text frame: {e}")
            pass
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        # default content box (will adjust if image present)
        content_left_in, content_top_in = page_margin, title_bottom_in
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        # Optional image placement with proper content reflow
        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                log.debug(f"Searching for image query: '{image_query}'")
                try:
                    log.debug(f"Downloading image from URL: {image_url}")
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)
                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0
                    log.debug(f"Image dimensions: {img_w_in} x {img_h_in} inches")

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = title_bottom_in
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)

                    slide.shapes.add_picture(image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                    log.debug(f"Image added at position: left={img_left_in}, top={img_top_in}")
                except Exception:
                    pass

        # apply content box geometry
        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        # estimate capacity to guide initial font size; autosize will fine-tune
        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        # Ensure positive dimensions to avoid calculation issues
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        # Ensure we still have a valid text frame after geometry changes
        try:
            tf = content_shape.text_frame
        except Exception:
            # If text frame access fails, try to recreate
            try:
                tf = content_shape.text_frame
            except Exception:
                log.warning("Could not access text frame for content shape")
                continue

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line) if line is not None else ""
            run.font.size = font_size
            p.space_after = PptPt(6)

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_word(content: list[dict] | str, filename: str, folder_path: str | None = None, title: str | None = None) -> dict:
    log.debug("Creating Word document")

    if isinstance(content, str):
        content = _convert_markdown_to_structured(content)
    elif not isinstance(content, list):
        content = []

    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "docx")

    # Template handling similar to create_presentation
    use_template = False
    doc = None

    if DOCX_TEMPLATE:
        try:
            src = DOCX_TEMPLATE
            if hasattr(DOCX_TEMPLATE, "paragraphs") and hasattr(DOCX_TEMPLATE, "save"):
                # If DOCX_TEMPLATE is already a Document object, create a copy
                buf = BytesIO()
                DOCX_TEMPLATE.save(buf)
                buf.seek(0)
                src = buf

            # Load template document
            doc = Document(src)
            use_template = True
            log.debug("Using DOCX template")

            # Properly clear existing content while preserving styles
            # Remove all paragraphs and tables
            for element in doc.element.body:
                if element.tag.endswith('}p') or element.tag.endswith('}tbl'):
                    doc.element.body.remove(element)

        except Exception as e:
            log.warning(f"Failed to load DOCX template: {e}")
            use_template = False
            doc = None

    if not use_template:
        doc = Document()
        log.debug("Creating new Word document without template")

    # Add title if provided
    if title:
        title_paragraph = doc.add_paragraph(title)
        try:
            title_paragraph.style = doc.styles['Title']
        except KeyError:
            # Fallback if template doesn't have Title style
            try:
                title_paragraph.style = doc.styles['Heading 1']
            except KeyError:
                # Manual formatting if no built-in styles available
                run = title_paragraph.runs[0] if title_paragraph.runs else title_paragraph.add_run()
                run.font.size = DocxPt(20)
                run.font.bold = True
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        log.debug("Document title added")

    # Add content to document
    for item in content or []:
        if isinstance(item, str):
            doc.add_paragraph(item)
        elif isinstance(item, dict):
            if item.get("type") == "image_query":
                new_item = {
                    "type": "image",
                    "query": item.get("query")
                }
                image_query = new_item.get("query")
                if image_query:
                    log.debug(f"Image search for the query : {image_query}")
                    image_url = search_image(image_query)
                    if image_url:
                        response = requests.get(image_url)
                        image_data = BytesIO(response.content)
                        doc.add_picture(image_data, width=Inches(6))
                        log.debug("Image successfully added")
                    else:
                        log.warning(f"Image search for : '{image_query}'")
            elif "type" in item:
                item_type = item.get("type")
                if item_type == "title":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 1']
                    except KeyError:
                        # Fallback if template doesn't have Heading 1 style
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(18)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Title added")
                elif item_type == "subtitle":
                    paragraph = doc.add_paragraph(item.get("text", ""))
                    try:
                        paragraph.style = doc.styles['Heading 2']
                    except KeyError:
                        # Fallback if template doesn't have Heading 2 style
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.font.size = DocxPt(16)
                        run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    log.debug("Subtitle added")
                elif item_type == "paragraph":
                    doc.add_paragraph(item.get("text", ""))
                    log.debug("Paragraph added")
                elif item_type == "list":
                    items = item.get("items", [])
                    for i, item_text in enumerate(items):
                        paragraph = doc.add_paragraph(item_text)
                        try:
                            paragraph.style = doc.styles['List Bullet']
                        except KeyError:
                            # Fallback if template doesn't have List Bullet style
                            paragraph.style = doc.styles['Normal']
                    log.debug("List added")
                elif item_type == "image":
                    image_query = item.get("query")
                    if image_query:
                        log.debug(f"Image search for the query : {image_query}")
                        image_url = search_image(image_query)
                        if image_url:
                            response = requests.get(image_url)
                            image_data = BytesIO(response.content)
                            doc.add_picture(image_data, width=Inches(6))
                            log.debug("Image successfully added")
                        else:
                            log.warning(f"Image search for : '{image_query}'")
                elif item_type == "table":
                    data = item.get("data", [])
                    if data:
                        # Check if template has existing tables to copy style from
                        template_table_style = None
                        if use_template and DOCX_TEMPLATE:
                            try:
                                # Look for existing tables in the original template
                                for table in DOCX_TEMPLATE.tables:
                                    if table.style:
                                        template_table_style = table.style
                                        break
                            except Exception:
                                pass
                        
                        table = doc.add_table(rows=len(data), cols=len(data[0]) if data else 0)
                        
                        # Apply template table style
                        if template_table_style:
                            try:
                                table.style = template_table_style
                                log.debug(f"Applied template table style: {template_table_style.name}")
                            except Exception as e:
                                log.debug(f"Could not apply template table style: {e}")
                        else:
                            # Try to apply a built-in table style
                            try:
                                # Try common built-in styles
                                for style_name in ['Table Grid', 'Light Grid Accent 1', 'Medium Grid 1 Accent 1', 'Light List Accent 1']:
                                    try:
                                        table.style = doc.styles[style_name]
                                        log.debug(f"Applied built-in table style: {style_name}")
                                        break
                                    except KeyError:
                                        continue
                            except Exception as e:
                                log.debug(f"Could not apply any table style: {e}")
                        
                        # Fill table data
                        for i, row in enumerate(data):
                            for j, cell in enumerate(row):
                                cell_obj = table.cell(i, j)
                                cell_obj.text = str(cell)
                                
                                # Apply header formatting to first row
                                if i == 0:
                                    for paragraph in cell_obj.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                        # Center align header text
                                        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        # Additional table formatting if no style was applied
                        if not template_table_style:
                            try:
                                
                                # Add borders to table
                                tbl = table._tbl
                                tblPr = tbl.tblPr
                                tblBorders = parse_xml(r'<w:tblBorders {}><w:top w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:left w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:right w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="000000"/><w:insideV w:val="single" w:sz="4" w:space="0" w:color="000000"/></w:tblBorders>'.format(nsdecls('w')))
                                tblPr.append(tblBorders)
                            except Exception as e:
                                log.debug(f"Could not add table borders: {e}")
                        
                        log.debug("Table added with improved styling")
            elif "text" in item:
                doc.add_paragraph(item["text"])
                log.debug("Paragraph added")
    
    doc.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}

def _create_raw_file(content: str, filename: str | None, folder_path: str | None = None) -> dict:
    log.debug("Creating raw file")
    if folder_path is None:
        folder_path = _generate_unique_folder()

    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "txt")

    if fname.lower().endswith(".xml") and isinstance(content, str) and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'

    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content or "")

    return {"url": _public_url(folder_path, fname), "path": filepath}

@mcp.tool()
def create_file(data: dict, persistent: bool = PERSISTENT_FILES) -> dict:
    """
    Use:
    data = {
      "format": "pdf|docx|pptx|xlsx|csv|txt|xml|py|cs|etc",
      "filename": "name.ext",
      "content": ...,
      "slides_data": [...],
      "title": "Optional title"
    }
    """
    log.debug("Creating file via tool")
    folder_path = _generate_unique_folder()
    format_type = (data.get("format") or "").lower()
    filename = data.get("filename")
    content = data.get("content")
    title = data.get("title")

    if format_type == "pdf":
        result = _create_pdf(content if isinstance(content, list) else [str(content or "")], filename, folder_path=folder_path)
    elif format_type == "pptx":
        result = _create_presentation(data.get("slides_data", []), filename, folder_path=folder_path, title=title)
    elif format_type == "docx":
        result = _create_word(content if content is not None else [], filename, folder_path=folder_path, title=title)
    elif format_type == "xlsx":
        result = _create_excel(content if content is not None else [], filename, folder_path=folder_path)
    elif format_type == "csv":
        result = _create_csv(content if content is not None else [], filename, folder_path=folder_path)
    else:
        use_filename = filename or f"export.{format_type or 'txt'}"
        result = _create_raw_file(content if content is not None else "", use_filename, folder_path=folder_path)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}

@mcp.tool()
def generate_and_archive(files_data: list[dict], archive_format: str = "zip", archive_name: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    """
    files_data = [
      {"format":"pdf","filename":"r1.pdf","content":[...],"title":"..."},
      {"format":"pptx","filename":"slides.pptx","slides_data":[...],"title":"..."},
      ...
    ]
    """
    log.debug("Generating archive via tool")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")

        try:
            if fmt == "pdf":
                res = _create_pdf(content if isinstance(content, list) else [str(content or "")], fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = _create_presentation(file_info.get("slides_data", []), fname, folder_path=folder_path, title=title)
            elif fmt == "docx":
                res = _create_word(content if content is not None else [], fname, folder_path=folder_path, title=title)
            elif fmt == "xlsx":
                res = _create_excel(content if content is not None else [], fname, folder_path=folder_path)
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise

        generated_paths.append(res["path"])

    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    archive_filename = f"{archive_basename}.zip" if archive_format.lower() not in ("7z", "tar.gz") else f"{archive_basename}.{archive_format}"
    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode='w') as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, 'w') as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}

if __name__ == "__main__":
    mcp.run()
